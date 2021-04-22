#!/usr/bin/env python
# coding: utf-8

# This is a healper module to add images to the existing powerpoint slides. 
# 



# In[8]:


# Imports
from pptx import Presentation
from pptx.util import Inches
import os
import random
from PIL import Image
import argparse


# In[10]:


def open_pp_presentation(inputfile,outputfile,img):
    prs = Presentation(inputfile)
    print(f'The count of slides in the input file is: {len(prs.slides)}')
    for i in range(1, len(prs.slides)):
          print(f'Currently working on slide {i} of presentation {inputfile} ')
          slide = prs.slides[i]
          shapes = slide.shapes
          print(f'slide shape is {shapes}')
          left = Inches(random.randint(1,5))
          top = Inches(random.randint(1,5))
          pic = slide.shapes.add_picture(img,left,top)
          
    prs.save(outputfile)
    
    


# In[11]:


def multiple_file_update(logo):
    for file in os.listdir(os.getcwd()):
        if file.endswith('.PPTX') or file.endswith('.pptx') or file.endswith('.PPT') or file.endswith('.ppt'):
            open_pp_presentation(file,file,logo)


# In[ ]:


def resize_image(logo_input,width, height):
    img = Image.open(logo_input)
    img = img.resize((width,height),Image.LANCZOS)
    img.save(logo_input)


# In[12]:


if __name__ == '__main__':
    logo_input = input(' Please type the name of the image with extension which needs to be superimposed in the slide')
    resize_ind = input('Do you want to resize the image before adding ? 1-Yes,2-No')
    if int(resize_ind) == 1:
        w = input('Please enter the width of image in pixels')
        h = input('Please enter the height of image in pixels')
        resize_image(logo_input, int(w), int(h))
    multiple_file_update(logo_input)


# In[ ]:




